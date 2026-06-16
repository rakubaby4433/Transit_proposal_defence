"""
Proposal Defense Slide Generator
=================================
Intelligent Public Transit Route Recommendation System for Kathmandu Valley
Group: THA080BCT006 / 012 / 022 / 030  |  Supervisor: Er. Anup Shrestha

Builds a .pptx that follows the IOE Thapathali "Slide Design Guidelines"
(Dinesh Baniya Kshatri, 2020) and the 12-minute Proposal Defense order.

Updated to reflect the A* search proposal:
  * Core algorithm: A* search with admissible Haversine heuristic
  * Composite edge cost: w = alpha*d + beta*t + gamma*k
  * Budget: NPR 2,200
  * A* reference (Hart, Nilsson & Raphael, 1968) added

Design rules enforced:
  * White background, black text on every slide
  * No animations
  * Cover: Title Arial Bold 44 | Members Arial Bold 22 | Dept/Campus/Date Arial 20 | NO footer
  * Titles: Arial Bold, centered, 36-40 pt
  * Body: Arial, left-justified, 28 pt main bullet / 24 pt sub-bullet
  * <= 6 main bullets, <= 8 words per bullet, <= 2-3 sub-bullets
  * Footer (all non-cover slides): date bottom-left, slide number bottom-right
  * One figure / table per slide

USAGE:
  pip install -r requirements.txt
  python generate_slides.py
  -> produces Proposal_Defense.pptx

OPTIONAL IMAGES (drop in same folder, auto-embedded if present):
  block_diagram.png   flowchart.png   gantt.png
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ----------------------------------------------------------------------------
# CONSTANTS
# ----------------------------------------------------------------------------
FONT = "Arial"
BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY_LINE = RGBColor(0x00, 0x00, 0x00)
HEADER_FILL = RGBColor(0xD9, 0xD9, 0xD9)

PRESENTATION_DATE = "June 16, 2026"
OUTPUT_FILE = "Proposal_Defense.pptx"

# 4:3 canvas (10 x 7.5 in) keeps Arial sizing identical to the guideline deck.
SLIDE_W = Inches(10)
SLIDE_H = Inches(7.5)

# ----------------------------------------------------------------------------
# HELPERS
# ----------------------------------------------------------------------------

def _set_white_background(slide):
    """Force a solid white fill on the slide background."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


def _blank_slide(prs):
    """Return a brand-new blank slide with white background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 6 = truly blank
    _set_white_background(slide)
    return slide


def _add_title(slide, text, size=38):
    """Centered Arial Bold title, 36-40 pt."""
    size = max(36, min(40, size))
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9.0), Inches(1.0))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    f = run.font
    f.name = FONT
    f.size = Pt(size)
    f.bold = True
    f.color.rgb = BLACK
    return box


def _add_footer(slide, page_no):
    """Date bottom-left, slide number bottom-right (non-cover slides only)."""
    # Date (left)
    left = slide.shapes.add_textbox(Inches(0.3), Inches(7.05), Inches(4.0), Inches(0.35))
    p = left.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = PRESENTATION_DATE
    r.font.name = FONT
    r.font.size = Pt(12)
    r.font.color.rgb = BLACK
    # Slide number (right)
    right = slide.shapes.add_textbox(Inches(8.7), Inches(7.05), Inches(1.0), Inches(0.35))
    p = right.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = str(page_no)
    r.font.name = FONT
    r.font.size = Pt(12)
    r.font.color.rgb = BLACK


def _add_bullets(slide, items, top=Inches(1.6), height=Inches(5.0),
                 left=Inches(0.8), width=Inches(8.6)):
    """
    Add a left-justified bullet list.
    `items` is a list of (text, level) tuples.
      level 0 -> main bullet, Arial 28, solid disc
      level 1 -> sub-bullet,  Arial 24, dash
    """
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for text, level in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        p.level = level
        p.space_after = Pt(10)
        run = p.add_run()
        bullet_char = "\u2022  " if level == 0 else "\u2013  "
        run.text = bullet_char + text
        f = run.font
        f.name = FONT
        f.size = Pt(28) if level == 0 else Pt(24)
        f.bold = False
        f.color.rgb = BLACK
        if level == 1:
            # indent sub-bullets
            pPr = p._pPr if p._pPr is not None else p.get_or_add_pPr()
            pPr.set('marL', str(Emu(Inches(0.5)).emu if hasattr(Emu(Inches(0.5)),'emu') else 457200))
    return box


def _style_cell(cell, text, bold=False, size=24, align=PP_ALIGN.LEFT, fill=None):
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    else:
        cell.fill.solid()
        cell.fill.fore_color.rgb = WHITE
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = BLACK


def _maybe_image_or_placeholder(slide, filename, caption,
                                left=Inches(1.3), top=Inches(1.7),
                                width=Inches(7.4), height=Inches(4.6)):
    """Embed image if the file exists, else draw a labeled placeholder box."""
    if os.path.exists(filename):
        slide.shapes.add_picture(filename, left, top, width=width)
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = BLACK
        shape.line.width = Pt(1.25)
        tf = shape.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = caption
        r.font.name = FONT
        r.font.size = Pt(24)
        r.font.bold = False
        r.font.color.rgb = BLACK


# ----------------------------------------------------------------------------
# BUILD PRESENTATION
# ----------------------------------------------------------------------------

def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    page = 0  # footer counter (starts after cover)

    # ---------------- SLIDE 1 : COVER (no footer) ----------------
    s = _blank_slide(prs)

    # Title (Arial Bold 44)
    box = s.shapes.add_textbox(Inches(0.5), Inches(0.7), Inches(9.0), Inches(1.8))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "An Intelligent Public Transit Route Recommendation System for Kathmandu Valley"
    r.font.name = FONT; r.font.size = Pt(44); r.font.bold = True; r.font.color.rgb = BLACK

    # Members (Arial Bold 22)
    members = [
        "Arun Raj Neupane (THA080BCT006)",
        "Bipin Tharu (THA080BCT012)",
        "Naresh Rasaili (THA080BCT022)",
        "Rakesh Kumar Yadav (THA080BCT030)",
    ]
    mbox = s.shapes.add_textbox(Inches(0.5), Inches(3.0), Inches(9.0), Inches(2.0))
    mtf = mbox.text_frame; mtf.word_wrap = True
    for i, m in enumerate(members):
        p = mtf.paragraphs[0] if i == 0 else mtf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = m
        r.font.name = FONT; r.font.size = Pt(22); r.font.bold = True; r.font.color.rgb = BLACK

    # Supervisor + Dept + Campus + Date (Arial 20)
    info = [
        "Project Supervisor: Er. Anup Shrestha",
        "Department of Electronics and Computer Engineering",
        "Institute of Engineering, Thapathali Campus",
        "Kathmandu, Nepal",
        PRESENTATION_DATE,
    ]
    ibox = s.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(9.0), Inches(2.0))
    itf = ibox.text_frame; itf.word_wrap = True
    for i, line in enumerate(info):
        p = itf.paragraphs[0] if i == 0 else itf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = line
        r.font.name = FONT; r.font.size = Pt(20); r.font.bold = False; r.font.color.rgb = BLACK

    # ---------------- SLIDE 2 : PRESENTATION OUTLINE ----------------
    page += 1
    s = _blank_slide(prs)
    _add_title(s, "Presentation Outline")
    _add_bullets(s, [
        ("Introduction", 0),
        ("Motivation", 0),
        ("Objectives and Scope", 0),
        ("Proposed Methodology", 0),
        ("Expected Results and Applications", 0),
        ("Timeline, Budget and References", 0),
    ])
    _add_footer(s, page)

    # ---------------- SLIDE 3 : INTRODUCTION ----------------
    page += 1
    s = _blank_slide(prs)
    _add_title(s, "Introduction")
    _add_bullets(s, [
        ("Public buses are heavily used in the valley", 0),
        ("However, route information remains informal and fragmented", 0),
        ("Routes are rarely available in digital form", 0),
        ("Stops can be modeled as graph nodes", 0),
        ("Route connections become weighted graph edges", 0),
        ("Finding a journey becomes a routing problem", 0),
    ])
    _add_footer(s, page)

    # ---------------- SLIDE 4 : MOTIVATION ----------------
    page += 1
    s = _blank_slide(prs)
    _add_title(s, "Motivation")
    _add_bullets(s, [
        ("Newcomers face real difficulty using buses", 0),
        ("Identifying the correct boarding point is hard", 0),
        ("Knowing whether a transfer is required", 0),
        ("Fare and travel time remain uncertain", 0),
        ("Distance alone ignores transfers and walking", 0),
        ("This motivates a multi-criteria solution", 0),
    ])
    _add_footer(s, page)

    # ---------------- SLIDE 5 : OBJECTIVES ----------------
    page += 1
    s = _blank_slide(prs)
    _add_title(s, "Objectives")
    _add_bullets(s, [
        ("To develop a web-based recommendation system", 0),
        ("To model the bus network as a weighted graph", 0),
        ("To compute optimal routes using A* search", 0),
        ("using an admissible Haversine heuristic", 1),
        ("using a composite distance-transfer-walking cost", 1),
        ("To show boarding, transfer and alighting points", 0),
    ])
    _add_footer(s, page)

    # ---------------- SLIDE 6 : SCOPE OF PROJECT ----------------
    page += 1
    s = _blank_slide(prs)
    _add_title(s, "Scope of Project")
    _add_bullets(s, [
        ("Covers Kathmandu, Lalitpur and Bhaktapur", 0),
        ("Focuses on selected major bus routes", 0),
        ("Uses a structured static route dataset", 0),
        ("No real-time traffic or GPS tracking", 0),
        ("Fare and travel time are approximate", 0),
        ("Serves commuters, students and tourists", 0),
    ])
    _add_footer(s, page)

    # ---------------- SLIDE 7 : METHODOLOGY - SYSTEM ARCHITECTURE ----------------
    page += 1
    s = _blank_slide(prs)
    _add_title(s, "Methodology: System Architecture")
    _maybe_image_or_placeholder(
        s, "block_diagram.png",
        "[ Insert Block Diagram (Fig 4-1) here ]\n\nUser  ->  Frontend (React + Leaflet)\n->  Backend API (FastAPI)\n->  Nearest-Stop Lookup  ->  Graph Builder  ->  A* Search\n(Route Dataset feeds lookup and graph)")
    _add_footer(s, page)

    # ---------------- SLIDE 8 : METHODOLOGY - GRAPH MODEL & COST ----------------
    page += 1
    s = _blank_slide(prs)
    _add_title(s, "Methodology: Graph Model & Cost")
    _add_bullets(s, [
        ("Each bus stop is a graph node", 0),
        ("Each connection is a weighted edge", 0),
        ("Edge cost:  w = alpha*d + beta*t + gamma*k", 0),
        ("d = distance, t = transfer, k = walking", 1),
        ("Higher beta avoids transfers", 1),
        ("Higher gamma reduces walking distance", 1),
    ])
    _add_footer(s, page)

    # ---------------- SLIDE 9 : METHODOLOGY - A* SEARCH ----------------
    page += 1
    s = _blank_slide(prs)
    _add_title(s, "Methodology: A* Search Algorithm")
    _add_bullets(s, [
        ("A* evaluates nodes using  f(n) = g(n) + h(n)", 0),
        ("g(n) is the actual cost from start", 1),
        ("h(n) is the Haversine distance to goal", 1),
        ("The heuristic never overestimates, so it is admissible", 0),
        ("This guarantees an optimal route", 0),
        ("A* expands fewer nodes than Dijkstra", 0),
    ])
    _add_footer(s, page)

    # ---------------- SLIDE 10 : METHODOLOGY - ALGORITHM FLOWCHART ----------------
    page += 1
    s = _blank_slide(prs)
    _add_title(s, "Methodology: Algorithm Flowchart")
    _maybe_image_or_placeholder(
        s, "flowchart.png",
        "[ Insert Flowchart (Fig 4-3) here ]\n\nQuery -> Validate -> Spatial Index + Haversine\n-> Build Weighted Graph -> A* Search\n-> Estimate Fare & Time -> Display on Map")
    _add_footer(s, page)

    # ---------------- SLIDE 11 : METHODOLOGY - INSTRUMENTATION ----------------
    page += 1
    s = _blank_slide(prs)
    _add_title(s, "Methodology: Instrumentation")
    _add_bullets(s, [
        ("Frontend using React.js and Leaflet.js", 0),
        ("Backend using Python and FastAPI", 0),
        ("Map services from OpenStreetMap", 0),
        ("Spatial index for nearest-stop lookup", 0),
        ("Dataset stored as JSON or SQLite", 0),
        ("Developed in Visual Studio Code", 0),
    ])
    _add_footer(s, page)

    # ---------------- SLIDE 12 : EXPECTED RESULTS (1) ----------------
    page += 1
    s = _blank_slide(prs)
    _add_title(s, "Expected Results")
    _add_bullets(s, [
        ("A structured Kathmandu Valley route dataset", 0),
        ("Optimal routing using A* search", 0),
        ("Recommendation of direct and transfer routes", 0),
        ("Clear boarding, transfer and alighting points", 0),
        ("Route shown on an interactive map", 0),
        ("Also presented as a textual step list", 0),
    ])
    _add_footer(s, page)

    # ---------------- SLIDE 13 : EXPECTED RESULTS (2) ----------------
    page += 1
    s = _blank_slide(prs)
    _add_title(s, "Expected Results (Contd.)")
    _add_bullets(s, [
        ("Approximate fare estimation", 0),
        ("Approximate travel time estimation", 0),
        ("A simple and user-friendly interface", 0),
        ("Efficient nearest-stop lookup via spatial index", 0),
        ("Improved travel experience for newcomers", 0),
        ("Centralized access to transit information", 0),
    ])
    _add_footer(s, page)

    # ---------------- SLIDE 14 : PROJECT APPLICATIONS ----------------
    page += 1
    s = _blank_slide(prs)
    _add_title(s, "Project Applications")
    _add_bullets(s, [
        ("Daily commuters within Kathmandu Valley", 0),
        ("Students travelling to their colleges", 0),
        ("Tourists exploring the valley", 0),
        ("Newcomers unfamiliar with bus routes", 0),
        ("Researchers studying urban mobility", 0),
        ("A foundation for future transit applications", 0),
    ])
    _add_footer(s, page)

    # ---------------- SLIDE 15 : TENTATIVE TIMELINE (GANTT) ----------------
    page += 1
    s = _blank_slide(prs)
    _add_title(s, "Tentative Timeline (Gantt Chart)")
    _maybe_image_or_placeholder(
        s, "gantt.png",
        "[ Insert Gantt Chart (Fig 6-1) here ]\n\nRequirement Study -> Data Collection -> Graph & Algorithm\n-> Frontend & Map Integration -> Testing -> Documentation")
    _add_footer(s, page)

    # ---------------- SLIDE 16 : ESTIMATED BUDGET (table) ----------------
    page += 1
    s = _blank_slide(prs)
    _add_title(s, "Estimated Project Budget")
    rows, cols = 4, 2
    tbl_w = Inches(7.0); tbl_h = Inches(2.6)
    left = Inches(1.5); top = Inches(2.0)
    table = s.shapes.add_table(rows, cols, left, top, tbl_w, tbl_h).table
    table.columns[0].width = Inches(4.6)
    table.columns[1].width = Inches(2.4)
    # header
    _style_cell(table.cell(0, 0), "Item", bold=True, size=24, align=PP_ALIGN.LEFT, fill=HEADER_FILL)
    _style_cell(table.cell(0, 1), "Estimated Cost (NPR)", bold=True, size=24, align=PP_ALIGN.CENTER, fill=HEADER_FILL)
    data = [
        ("Internet / data collection", "1,000"),
        ("Documentation", "700"),
        ("Miscellaneous", "500"),
    ]
    for i, (item, cost) in enumerate(data, start=1):
        _style_cell(table.cell(i, 0), item, bold=False, size=24, align=PP_ALIGN.LEFT)
        _style_cell(table.cell(i, 1), cost, bold=False, size=24, align=PP_ALIGN.CENTER)
    # total note
    total_box = s.shapes.add_textbox(Inches(1.5), Inches(4.9), Inches(7.0), Inches(0.6))
    p = total_box.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    r = p.add_run(); r.text = "Total:  NPR 2,200"
    r.font.name = FONT; r.font.size = Pt(28); r.font.bold = True; r.font.color.rgb = BLACK
    # note about open-source
    note_box = s.shapes.add_textbox(Inches(1.5), Inches(5.6), Inches(7.0), Inches(0.8))
    p = note_box.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = "All software and map tiles are free and open-source."
    r.font.name = FONT; r.font.size = Pt(20); r.font.bold = False; r.font.color.rgb = BLACK
    _add_footer(s, page)

    # ---------------- SLIDE 17 : REFERENCES (1) ----------------
    page += 1
    s = _blank_slide(prs)
    _add_title(s, "References")
    refs1 = [
        "[1] E. W. Dijkstra, \u201cA note on two problems in connexion with graphs,\u201d Numerische Mathematik, vol. 1, pp. 269-271, 1959.",
        "[2] P. E. Hart, N. J. Nilsson and B. Raphael, \u201cA formal basis for the heuristic determination of minimum cost paths,\u201d IEEE Trans. SSC, vol. 4, no. 2, pp. 100-107, 1968.",
        "[3] T. H. Cormen et al., Introduction to Algorithms, 3rd ed. MIT Press, 2009.",
        "[4] M. A. Rahman et al., \u201cShortest path algorithms in transportation networks: A survey,\u201d IJCA, vol. 181, no. 14, 2018.",
        "[5] S. Zhang et al., \u201cA survey of route recommendations,\u201d J. Transp. Tech., vol. 12, no. 3, 2022.",
    ]
    box = s.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(8.8), Inches(5.0))
    tf = box.text_frame; tf.word_wrap = True
    for i, ref in enumerate(refs1):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT; p.space_after = Pt(10)
        r = p.add_run(); r.text = ref
        r.font.name = FONT; r.font.size = Pt(18); r.font.color.rgb = BLACK
    _add_footer(s, page)

    # ---------------- SLIDE 18 : REFERENCES (2) ----------------
    page += 1
    s = _blank_slide(prs)
    _add_title(s, "References (Contd.)")
    refs2 = [
        "[6] MobilityData, GTFS Reference. [Online]. Available: https://gtfs.org/",
        "[7] OpenStreetMap. [Online]. Available: https://www.openstreetmap.org/",
        "[8] V. Agafonkin, Leaflet: An open-source JavaScript library. [Online]. Available: https://leafletjs.com/",
        "[9] Google Maps. [Online]. Available: https://www.google.com/maps",
        "[10] Moovit transit app. [Online]. Available: https://moovitapp.com/",
    ]
    box = s.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(8.8), Inches(5.0))
    tf = box.text_frame; tf.word_wrap = True
    for i, ref in enumerate(refs2):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT; p.space_after = Pt(10)
        r = p.add_run(); r.text = ref
        r.font.name = FONT; r.font.size = Pt(18); r.font.color.rgb = BLACK
    _add_footer(s, page)

    # ---------------- SLIDE 19 : THANK YOU ----------------
    page += 1
    s = _blank_slide(prs)
    box = s.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(9.0), Inches(1.5))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "Thank You!"
    r.font.name = FONT; r.font.size = Pt(40); r.font.bold = True; r.font.color.rgb = BLACK
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = "Questions & Discussion"
    r2.font.name = FONT; r2.font.size = Pt(28); r2.font.bold = False; r2.font.color.rgb = BLACK
    _add_footer(s, page)

    prs.save(OUTPUT_FILE)
    print(f"Saved {OUTPUT_FILE} with {len(prs.slides._sldIdLst)} slides.")


if __name__ == "__main__":
    build()
